#!/usr/bin/env python3
"""Build editable Phase 5 operations decks with existing approved brand assets."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output" / "phase5" / "presentations"
OUT.mkdir(parents=True, exist_ok=True)

BLACK = RGBColor(8, 8, 9)
CHARCOAL = RGBColor(28, 28, 30)
GOLD = RGBColor(217, 164, 65)
GOLD_LIGHT = RGBColor(246, 232, 191)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(180, 180, 184)
GREEN = RGBColor(68, 180, 105)
AMBER = RGBColor(226, 162, 58)
RED = RGBColor(214, 75, 75)

LOGO = ROOT / "public" / "brand" / "black-diamond" / "our-town-logo.png"
HERO = ROOT / "public" / "brand" / "black-diamond" / "hero-home-desktop.jpg"
HERO_MOBILE = ROOT / "public" / "brand" / "black-diamond" / "hero-home-mobile.jpg"
SOCIAL = ROOT / "public" / "brand" / "black-diamond" / "hero-social-4x5.jpg"
HEADSHOT = ROOT / "public" / "images" / "ask-magic-mike" / "brand-pack-v2" / "mike-headshot-source.jpg"
THINKING = ROOT / "output" / "phase5" / "build-assets" / "chat-widget-concept.png"
EXPLAINING = HEADSHOT
NOTIFICATION = ROOT / "public" / "images" / "ask-magic-mike" / "notifications" / "lead-alert-frame-v1.png"
BRAND_BOARD = ROOT / "public" / "images" / "ask-magic-mike" / "brand-pack-v2" / "brand-board-v2-web.jpg"


def blank_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    return prs


def add_background(slide, color=BLACK) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_logo(slide, x=0.62, y=0.36, w=1.8) -> None:
    slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), width=Inches(w))


def add_footer(slide, number: int, label="PHASE 5 LIVE OPERATIONS") -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(7.12), Inches(12.1), Inches(0.015))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()
    add_text(slide, label, 0.65, 7.18, 4.6, 0.18, 8.5, MUTED, bold=True)
    add_text(slide, f"{number:02d}", 12.15, 7.17, 0.5, 0.18, 8.5, GOLD, bold=True, align=PP_ALIGN.RIGHT)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: float, color=WHITE, *, bold=False, align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP, margin=0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, number: int, kicker="ASK MAGIC MIKE") -> None:
    add_logo(slide)
    add_text(slide, kicker, 9.35, 0.42, 3.35, 0.22, 9, GOLD, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, title, 0.65, 1.14, 12.0, 0.72, 36, WHITE, bold=True)
    add_footer(slide, number)


def add_notes(slide, narrative: str, sources: list[str]) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = narrative + "\n\n[Sources]\n" + "\n".join(f"- {source}" for source in sources)


def add_metric(slide, x: float, y: float, value: str, label: str, color=GOLD, width=2.5) -> None:
    add_text(slide, value, x, y, width, 0.64, 32, color, bold=True)
    add_text(slide, label.upper(), x, y + 0.66, width, 0.42, 11, MUTED, bold=True)


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size=18, color=WHITE) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True
    frame.margin_left = Inches(0); frame.margin_right = Inches(0)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = "•  " + item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
    return box


def add_status_row(
    slide,
    y: float,
    label: str,
    status: str,
    detail: str,
    color=GREEN,
    detail_width: float = 6.1,
    detail_size: float = 13,
) -> None:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.76), Inches(y + 0.08), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    add_text(slide, label, 1.06, y, 2.8, 0.3, 16, WHITE, bold=True)
    add_text(slide, status, 4.0, y, 2.05, 0.3, 14, color, bold=True)
    add_text(slide, detail, 6.2, y, detail_width, 0.48, detail_size, MUTED)


def add_photo_panel(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + h - 0.055), Inches(w), Inches(0.055))
    overlay.fill.solid(); overlay.fill.fore_color.rgb = GOLD; overlay.line.fill.background()


def make_phase5_deck() -> Path:
    prs = blank_deck()

    # 1 - cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_photo_panel(slide, HERO, 7.95, 0, 5.38, 7.5)
    add_logo(slide, 0.72, 0.55, 2.2)
    add_text(slide, "PHASE 5", 0.75, 2.0, 3.0, 0.28, 13, GOLD, bold=True)
    add_text(slide, "Live operations\nand controlled growth", 0.75, 2.42, 6.4, 1.9, 46, WHITE, bold=True)
    add_text(slide, "A clean queue, protected first-lead handling, and an evidence-first operating cadence.", 0.78, 4.72, 5.95, 0.82, 19, GOLD_LIGHT)
    add_text(slide, "Verified August 15, 2026", 0.78, 6.62, 3.8, 0.28, 11, MUTED)
    add_notes(slide, "Open with the operating result: the system is live and monitored, while zero genuine leads remains an honest current state.", ["docs/PHASE5_PRECHANGE_PRODUCTION_SNAPSHOT.md", "docs/PHASE5_RELEASE_NOTES.md"])

    # 2
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "The system is live and the evidence is current", 2)
    add_metric(slide, 0.78, 2.2, "READY", "Vercel production", GREEN, 2.4)
    add_metric(slide, 3.3, 2.2, "9 / 9", "Public monitor", GREEN, 2.1)
    add_metric(slide, 5.72, 2.2, "0", "Three-hour runtime errors", GREEN, 2.4)
    add_metric(slide, 8.45, 2.2, "1", "Canonical database", GOLD, 2.2)
    add_metric(slide, 10.75, 2.2, "1", "Canonical form", GOLD, 1.8)
    add_text(slide, "Production commit", 0.8, 4.15, 2.2, 0.3, 13, MUTED, bold=True)
    add_text(slide, "e754456cecaf6538df25bb4bf5eebe57ebf6eacb", 0.8, 4.52, 7.4, 0.35, 17, WHITE, bold=True)
    add_text(slide, "Deployment", 0.8, 5.22, 2.2, 0.3, 13, MUTED, bold=True)
    add_text(slide, "dpl_3ogimm1EhHCaPkEfXLAeojrm2H8Z", 0.8, 5.58, 6.0, 0.35, 17, WHITE, bold=True)
    add_text(slide, "Neon production", 8.5, 4.15, 2.4, 0.3, 13, MUTED, bold=True)
    add_text(slide, "bitter-star-20214385\nbr-round-base-auh6h2wd", 8.5, 4.52, 3.7, 0.78, 17, WHITE, bold=True)
    add_notes(slide, "State only evidence already verified against authenticated production accounts and public health checks.", ["docs/PHASE5_PRECHANGE_PRODUCTION_SNAPSHOT.md"])

    # 3
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "A clean live queue is now a release invariant", 3)
    chart_data = ChartData(); chart_data.categories = ["Genuine live", "Suppressed QA", "QA in Active/New"]
    chart_data.add_series("Records", (0, 6, 0))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.75), Inches(2.05), Inches(6.3), Inches(4.35), chart_data).chart
    chart.has_legend = False; chart.has_title = False
    chart.value_axis.minimum_scale = 0; chart.value_axis.maximum_scale = 7; chart.value_axis.major_unit = 1
    chart.value_axis.tick_labels.font.color.rgb = MUTED; chart.category_axis.tick_labels.font.color.rgb = WHITE
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = GOLD
    add_bullets(slide, ["TEST badge remains visible", "Creation timestamps remain truthful", "No QA record can be STALLED or ROUTING READY", "Test and suppressed records are excluded from live reporting", "Explicit evidence is required before classifying a lead as QA"], 7.55, 2.18, 4.95, 3.8, 18)
    add_notes(slide, "Explain that the six QA records are intentionally retained evidence, not deleted records and not sales opportunities.", ["docs/LEAD_QUEUE_INVARIANTS.md", "docs/PHASE5_PRECHANGE_PRODUCTION_SNAPSHOT.md"])

    # 4
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "The signed-in audit closed misleading states", 4)
    add_photo_panel(slide, BRAND_BOARD, 8.65, 1.92, 3.95, 4.65)
    add_status_row(slide, 2.18, "Active/New", "CLEAN", "Zero QA or suppressed records", GREEN)
    add_status_row(slide, 2.96, "Closed/Test", "PRESERVED", "Six TEST-badged records with valid timestamps", GOLD)
    add_status_row(slide, 3.74, "Stalled", "CLEAR", "No QA-derived escalation state", GREEN)
    add_status_row(slide, 4.52, "Routing ready", "CLEAR", "No QA-derived routing signal", GREEN)
    add_status_row(slide, 5.30, "Live reporting", "PROTECTED", "Defensive live-ID filtering", GREEN)
    add_notes(slide, "Describe the signed-in evidence without displaying private lead details or treating the design board as production proof.", ["docs/PHASE5_PRECHANGE_PRODUCTION_SNAPSHOT.md", "tests/adminops/neon-reporting-exclusions.test.ts"])

    # 5
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "Brandon can operate the Lead Center now", 5)
    add_photo_panel(slide, THINKING, 8.72, 1.92, 3.75, 4.9)
    add_bullets(slide, ["Administrator role verified", "Email ownership verified", "Password credential present", "Authenticated Lead Center access confirmed", "Can review queue, assignments, notification state, and audit history", "Web Push remains a physical browser-permission step"], 0.85, 2.1, 7.15, 4.55, 20)
    add_notes(slide, "This slide is operational enablement, not an identity or credential disclosure. Do not show reset links, cookies, or session tokens.", ["docs/BRANDON_OPERATOR_ACTIVATION_STATUS.md", "docs/BRANDON_WEB_PUSH_ENROLLMENT.md"])

    # 6
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "Mike is provisioned; one private step remains", 6)
    add_photo_panel(slide, HEADSHOT, 8.72, 1.92, 3.75, 4.9)
    add_metric(slide, 0.85, 2.18, "PRIMARY", "Lead owner role", GOLD, 2.6)
    add_metric(slide, 3.6, 2.18, "0", "Active sessions", WHITE, 2.2)
    add_metric(slide, 5.85, 2.18, "PENDING", "Private password", AMBER, 2.4)
    add_bullets(slide, ["No password was chosen for Mike", "No reset token appears in shared evidence", "Email notifications continue independently", "After private activation, Mike can enroll one primary Push device"], 0.9, 4.15, 6.95, 2.0, 18)
    add_notes(slide, "The exact next action is Mike choosing his own password through a valid private activation flow.", ["docs/MIKE_PRIMARY_OWNER_ACTIVATION_STATUS.md", "docs/MIKE_WEB_PUSH_ENROLLMENT.md"])

    # 7
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "The first real lead follows one protected path", 7)
    labels = [("1", "STORE", "Durable Neon record"), ("2", "CLASSIFY", "Explicit live/test evidence"), ("3", "SCORE", "Deterministic factors"), ("4", "ASSIGN", "Approved owner + reason"), ("5", "NOTIFY", "Email + optional Push"), ("6", "RESPOND", "Consent-permitted human action")]
    x_positions = [0.72, 2.82, 4.92, 7.02, 9.12, 11.22]
    for idx in range(len(labels) - 1):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_positions[idx] + 1.45), Inches(3.55), Inches(0.65), Inches(0.035))
        line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()
    for x, (num, label, detail) in zip(x_positions, labels):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.78), Inches(1.45), Inches(1.45))
        circle.fill.solid(); circle.fill.fore_color.rgb = CHARCOAL; circle.line.color.rgb = GOLD; circle.line.width = Pt(2)
        add_text(slide, num, x, 3.08, 1.45, 0.44, 28, GOLD, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x - 0.2, 4.48, 1.85, 0.3, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, x - 0.25, 4.9, 1.95, 0.78, 11, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "No fabricated prospect  •  no carrier SMS  •  no public response-time promise", 1.25, 6.24, 10.9, 0.38, 16, GOLD_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Walk left to right. Storage precedes every notification; QA classification requires explicit evidence; customer action follows consent.", ["docs/FIRST_LIVE_LEAD_RESPONSE_RUNBOOK.md", "docs/LEAD_QUEUE_INVARIANTS.md", "src/lib/operations/first-live-lead-monitor.ts"])

    # 8
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "Form 3 stays live; held forms advance one at a time", 8)
    form_rows = [("1", "Contact", "HELD", AMBER), ("2", "Seller options", "HELD", AMBER), ("3", "Home value", "ACTIVE", GREEN), ("4", "Recruiting", "SEPARATE", MUTED), ("5", "Rental search", "HELD", AMBER), ("6", "Short-term rental", "HELD", AMBER), ("7", "Property alerts", "LEGAL HOLD", RED)]
    for idx, (fid, purpose, status, color) in enumerate(form_rows):
        y = 2.08 + idx * 0.62
        add_text(slide, f"FORM {fid}", 0.82, y, 1.3, 0.28, 12, GOLD, bold=True)
        add_text(slide, purpose, 2.22, y, 3.25, 0.3, 16, WHITE, bold=True)
        add_text(slide, status, 5.72, y, 1.9, 0.3, 13, color, bold=True)
        detail = "Signed bridge; native duplicate alert inactive" if fid == "3" else ("Protected entry 1550; no retroactive subscription" if fid == "7" else "No canonical activation without BIC/owner gate")
        add_text(slide, detail, 7.5, y, 5.0, 0.38, 12.5, MUTED)
    add_notes(slide, "Form 4 is not a consumer lead form. Entry 1550 remains protected and is never repurposed for QA.", ["docs/GRAVITY_FORMS_BIC_APPROVAL_PACKET.md", "docs/PHASE5_PRECHANGE_PRODUCTION_SNAPSHOT.md"])

    # 9
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "Email is primary; Push follows physical enrollment", 9)
    add_photo_panel(slide, NOTIFICATION, 8.75, 1.96, 3.68, 4.78)
    add_status_row(slide, 2.18, "Internal email", "WIRED", "Durable outbox + bounded retry", GREEN, 2.25, 11.5)
    add_status_row(slide, 3.06, "Audit BCC", "PROTECTED", "Securely configured; never displayed", GREEN, 2.25, 11.5)
    add_status_row(slide, 3.94, "Brandon Push", "PENDING", "Physical browser permission", AMBER, 2.25, 11.5)
    add_status_row(slide, 4.82, "Mike Push", "PENDING", "Private login, then permission", AMBER, 2.25, 11.5)
    add_status_row(slide, 5.70, "Carrier SMS", "DISABLED", "No unsafe free-SMS workaround", MUTED, 2.25, 11.5)
    add_notes(slide, "The notification image is an approved visual template, not evidence of a genuine lead or delivered message.", ["docs/EMAIL_DELIVERY_SPEC.md", "docs/WEB_PUSH_DEVICE_REGISTER.csv", "public/images/ask-magic-mike/notifications/lead-alert-new-v2.png"])

    # 10
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "Monitors protect the live path without noise", 10)
    add_metric(slide, 0.85, 2.08, "2 MIN", "First-live reconciliation", GOLD, 2.5)
    add_metric(slide, 3.45, 2.08, "1 HR", "SLA sweep", GOLD, 2.2)
    add_metric(slide, 5.86, 2.08, "1 HR", "Public synthetic", GOLD, 2.2)
    add_metric(slide, 8.28, 2.08, "0", "Current live failures", GREEN, 2.2)
    add_metric(slide, 10.55, 2.08, "0", "TLS warnings", GREEN, 2.0)
    add_bullets(slide, ["Unassigned live leads escalate", "Failed live notifications remain visible and retry", "Unsuppressed QA makes first-live cron unhealthy", "Routine all-clear messages are suppressed", "Correlation IDs support private diagnosis"], 0.9, 4.25, 11.5, 2.0, 18)
    add_notes(slide, "Targets are internal operating controls, not consumer-facing guarantees.", ["vercel.json", "docs/INCIDENT_ESCALATION_MATRIX.csv", "src/lib/operations/first-live-lead-monitor.ts"])

    # 11
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "Owned traffic can grow without paid media", 11)
    add_photo_panel(slide, SOCIAL, 8.82, 1.96, 3.58, 4.82)
    add_bullets(slide, ["14 tagged placements staged", "UTM source, medium, campaign, and placement standardized", "QR package covers web, social, print, listing, rental, and open-house contexts", "Thirty-day content queue remains draft until BIC/content approval", "Paid media remains inactive"], 0.85, 2.08, 7.1, 3.8, 20)
    add_text(slide, "EXTERNAL PROMOTIONAL ASSETS PUBLISHED: 0", 0.88, 6.12, 6.7, 0.38, 16, GOLD_LIGHT, bold=True)
    add_notes(slide, "No social, GBP, email, or print asset is represented as published. All links remain staged until approval.", ["docs/ASK_MAGIC_MIKE_UTM_LINK_LIBRARY.csv", "docs/ZERO_SPEND_CONTENT_APPROVAL_QUEUE.csv"])

    # 12
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "Reports measure outcomes, not activity theater", 12)
    measures = [("LEADS", "Genuine only"), ("SPEED", "Assignment + contact"), ("QUALITY", "Score + qualification"), ("DELIVERY", "Provider outcomes"), ("SOURCE", "First + last touch"), ("VALUE", "Appointments to closings")]
    for idx, (label, detail) in enumerate(measures):
        x = 0.78 + (idx % 3) * 4.15
        y = 2.08 + (idx // 3) * 2.08
        add_text(slide, label, x, y, 3.45, 0.38, 23, GOLD, bold=True)
        add_text(slide, detail, x, y + 0.52, 3.45, 0.38, 17, WHITE)
    add_text(slide, "Day 1 establishes control. Day 7 shows operating reliability. Day 30 supports the next evidence-based decision.", 0.85, 6.13, 11.7, 0.42, 18, GOLD_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Business outcomes are recorded only when they occur. Assumptions remain visibly separate from evidence.", ["output/phase5/spreadsheets/ASK_MAGIC_MIKE_DAY1_OPERATIONS_REPORT.xlsx", "output/phase5/spreadsheets/ASK_MAGIC_MIKE_7_DAY_OPERATIONS_REPORT.xlsx", "output/phase5/spreadsheets/EVIDENCE_AND_ASSUMPTION_REGISTER.xlsx"])

    # 13
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, "One private action unlocks Mike's workflow", 13)
    add_photo_panel(slide, EXPLAINING, 8.8, 1.82, 3.65, 4.95)
    add_text(slide, "NEXT ACTION", 0.86, 2.1, 2.8, 0.34, 14, GOLD, bold=True)
    add_text(slide, "Mike chooses his private\nLead Center password.", 0.86, 2.62, 7.2, 1.28, 40, WHITE, bold=True)
    add_text(slide, "The funnel, email delivery, Form 3, monitoring, and Brandon's administrator access continue while that private step remains open.", 0.9, 4.46, 6.85, 1.1, 20, GOLD_LIGHT)
    add_text(slide, "No credential or activation token belongs in shared artifacts.", 0.9, 6.2, 6.4, 0.32, 14, MUTED)
    add_notes(slide, "Close with one precise, unavoidable human action and make clear that it does not block public lead capture.", ["docs/MIKE_PRIMARY_OWNER_ACTIVATION_STATUS.md", "docs/PHASE5_PRECHANGE_PRODUCTION_SNAPSHOT.md"])

    path = OUT / "ASK_MAGIC_MIKE_PHASE5_OPERATIONS_PRESENTATION.pptx"
    prs.save(path)
    return path


def make_30_day_deck() -> Path:
    prs = blank_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_photo_panel(slide, HERO_MOBILE, 8.6, 0, 4.73, 7.5); add_logo(slide, 0.72, 0.55, 2.2)
    add_text(slide, "30-DAY EXECUTIVE REPORT", 0.75, 2.0, 5.8, 0.3, 13, GOLD, bold=True)
    add_text(slide, "Evidence before\nexpansion", 0.75, 2.5, 7.0, 1.55, 46, WHITE, bold=True)
    add_text(slide, "Editable operating review for live leads, source quality, response, delivery, and next-month decisions.", 0.78, 4.38, 6.7, 1.0, 19, GOLD_LIGHT)
    add_text(slide, "Current baseline: Day 0 / no genuine lead fabricated", 0.78, 6.58, 5.4, 0.28, 11, MUTED)
    add_notes(slide, "This is an editable report framework. Populate actual outcomes after 30 live operating days.", ["docs/PHASE5_PRECHANGE_PRODUCTION_SNAPSHOT.md", "output/phase5/spreadsheets/ASK_MAGIC_MIKE_OPERATING_SCOREBOARD.xlsx"])
    slides = [
        ("Start with data integrity", ["Genuine leads only", "QA and suppressed records excluded", "Attribution reconciled", "No live customer PII in this deck"]),
        ("Show the funnel with actual outcomes", ["Lead created", "Qualified", "Appointment", "Signed client", "Closing"]),
        ("Compare source quality, not just volume", ["Domain and placement", "First and last touch", "Qualified rate", "Appointment rate", "Attributed revenue"]),
        ("Test whether response operations held", ["Assignment time", "First-contact time", "Unassigned count", "SLA breaches", "Human acknowledgment"]),
        ("Reconcile every delivery channel", ["Internal email delivered", "Retry outcomes", "Hidden audit-copy result", "Web Push enrolled devices", "Carrier SMS remains disabled"]),
        ("Review form expansion one gate at a time", ["Form 3 remains baseline", "Held forms require consent and routing approval", "One QA submission per approved form", "Rollback stays available"]),
        ("Make the next decision from observed evidence", ["Continue current owned placements", "Improve one proven conversion bottleneck", "Approve one held form or keep it held", "Keep paid media inactive until the live pipe proves reliable"]),
    ]
    for idx, (title, bullets) in enumerate(slides, 2):
        slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, title, idx, kicker="30-DAY EXECUTIVE REPORT")
        add_bullets(slide, bullets, 0.95, 2.25, 7.2, 3.9, 22)
        add_metric(slide, 9.1, 2.35, "0", "Current verified baseline", GOLD, 2.8)
        add_text(slide, "Replace with observed 30-day actuals. Never fill with synthetic performance.", 8.8, 4.02, 3.65, 1.15, 17, GOLD_LIGHT, bold=True, align=PP_ALIGN.CENTER)
        add_notes(slide, "Use only actual operating results from the canonical system and approved business records.", ["output/phase5/spreadsheets/ASK_MAGIC_MIKE_OPERATING_SCOREBOARD.xlsx", "output/phase5/spreadsheets/EVIDENCE_AND_ASSUMPTION_REGISTER.xlsx"])
    path = OUT / "ASK_MAGIC_MIKE_30_DAY_EXECUTIVE_REPORT.pptx"
    prs.save(path)
    return path


def main() -> None:
    paths = [make_phase5_deck(), make_30_day_deck()]
    print("\n".join(str(path) for path in paths))


if __name__ == "__main__":
    main()
